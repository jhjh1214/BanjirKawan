"""Build BanjirKawan.pptx from docs/pptx-build-spec.md.

Generates the 10-slide, 16:9 pitch deck with headlines, body copy, the system-map
and loop diagrams, metric tiles, the worst-day table, and verbatim speaker notes.
Video/screenshots are rendered as labelled placeholders unless files exist in
docs/assets/ (demo.mp4, map.png, report.png, metrics.png).

Run:  python scripts/build_deck.py
"""
from __future__ import annotations
import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE
from pptx.oxml.ns import qn

# ---- palette ----
INK          = RGBColor(0x0F, 0x17, 0x2A)
SECONDARY    = RGBColor(0x47, 0x55, 0x69)
ACCENT       = RGBColor(0x02, 0x84, 0xC7)
ACCENT_LIGHT = RGBColor(0x0E, 0xA5, 0xE9)
WHITE        = RGBColor(0xFF, 0xFF, 0xFF)
BG_SOFT      = RGBColor(0xF8, 0xFA, 0xFC)
CARD_BORDER  = RGBColor(0xE2, 0xE8, 0xF0)
GREEN        = RGBColor(0x16, 0xA3, 0x4A)
RED          = RGBColor(0xDC, 0x26, 0x26)
GREY_BOX     = RGBColor(0xE2, 0xE8, 0xF0)
LIGHT_BLUE   = RGBColor(0xE0, 0xF2, 0xFE)
LIGHT_GREEN  = RGBColor(0xDC, 0xFC, 0xE7)
LIGHT_GREY   = RGBColor(0xF1, 0xF5, 0xF9)
FONT = "Calibri"

HERE = os.path.dirname(os.path.abspath(__file__))
ROOT = os.path.dirname(HERE)
ASSETS = os.path.join(ROOT, "docs", "assets")

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]


# ---------- helpers ----------
def slide():
    return prs.slides.add_slide(BLANK)


def _no_line(shape):
    shape.line.fill.background()


def add_bg(s, color=WHITE):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid(); r.fill.fore_color.rgb = color; _no_line(r)
    r._element.addprevious(r._element)  # keep at back (added first anyway)
    return r


def add_text(s, l, t, w, h, text, size, color=INK, bold=False, italic=False,
             align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, font=FONT):
    tb = s.shapes.add_textbox(Inches(l), Inches(t), Inches(w), Inches(h))
    tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
    lines = text.split("\n")
    for i, line in enumerate(lines):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        r = p.add_run(); r.text = line
        f = r.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = font
    return tb


def add_headline(s, text):
    return add_text(s, 0.6, 0.45, 12.13, 0.9, text, 40, INK, bold=True)


def add_footer(s):
    add_text(s, 0.6, 7.05, 12.13, 0.35,
             "BanjirKawan · Climate Systems Hackathon 2026", 11, SECONDARY)


def add_bar(s, l, t, w, h=0.09, color=ACCENT):
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = color; _no_line(r)
    return r


def box(s, l, t, w, h, fill=BG_SOFT, border=CARD_BORDER, rounded=True, border_w=1.0):
    shp = MSO_SHAPE.ROUNDED_RECTANGLE if rounded else MSO_SHAPE.RECTANGLE
    r = s.shapes.add_shape(shp, Inches(l), Inches(t), Inches(w), Inches(h))
    r.fill.solid(); r.fill.fore_color.rgb = fill
    if border is None:
        _no_line(r)
    else:
        r.line.color.rgb = border; r.line.width = Pt(border_w)
    if rounded:
        try:
            r.adjustments[0] = 0.06
        except Exception:
            pass
    r.shadow.inherit = False
    return r


def box_text(s, l, t, w, h, text, size, color=INK, bold=False, fill=BG_SOFT,
             border=CARD_BORDER, align=PP_ALIGN.CENTER, rounded=True, italic=False):
    r = box(s, l, t, w, h, fill=fill, border=border, rounded=rounded)
    tf = r.text_frame; tf.word_wrap = True; tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    tf.margin_left = Inches(0.1); tf.margin_right = Inches(0.1)
    for i, line in enumerate(text.split("\n")):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        run = p.add_run(); run.text = line
        f = run.font; f.size = Pt(size); f.bold = bold; f.italic = italic
        f.color.rgb = color; f.name = FONT
    return r


def arrow(s, x1, y1, x2, y2, color=ACCENT, weight=2.0, dash=False):
    cn = s.shapes.add_connector(2, Inches(x1), Inches(y1), Inches(x2), Inches(y2))  # 2=straight
    ln = cn.line; ln.color.rgb = color; ln.width = Pt(weight)
    lnEl = ln._get_or_add_ln()
    tail = lnEl.makeelement(qn("a:tailEnd"), {"type": "triangle"})
    lnEl.append(tail)
    if dash:
        d = lnEl.makeelement(qn("a:prstDash"), {"val": "dash"})
        lnEl.append(d)
    return cn


def oval(s, l, t, w, h, color=RED, weight=2.5):
    o = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(l), Inches(t), Inches(w), Inches(h))
    o.fill.background(); o.line.color.rgb = color; o.line.width = Pt(weight)
    o.shadow.inherit = False
    return o


def notes(s, text):
    s.notes_slide.notes_text_frame.text = text


def asset(name):
    p = os.path.join(ASSETS, name)
    return p if os.path.exists(p) else None


# ============ SLIDE 1 ============
s = slide(); add_bg(s, BG_SOFT)
add_text(s, 0.6, 2.6, 12.13, 1.2, "BanjirKawan", 60, ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.95, 10.33, 0.9,
         "Turning flood warnings into shop-level survival plans — and survival plans into claim evidence.",
         22, SECONDARY, italic=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 5.0, 10.33, 0.5,
         "Problem 1  ·  Team of 3 — [Name], [Name], [Name]", 15, SECONDARY, align=PP_ALIGN.CENTER)
notes(s, "(No speech on the title — go straight into the hook on Slide 2.)")

# ============ SLIDE 2 ============
s = slide(); add_bg(s)
add_headline(s, "The warning existed. It never reached her shop floor.")
c = box(s, 0.6, 1.7, 5.6, 4.4, fill=BG_SOFT)
add_text(s, 1.0, 2.15, 4.8, 0.5, "Kak Ros", 26, INK, bold=True)
add_text(s, 1.0, 2.7, 4.8, 0.4, "kedai runcit · Taman Sri Muda", 18, SECONDARY)
add_text(s, 1.0, 3.5, 4.8, 1.1, "−RM 40,000", 54, RED, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 4.75, 4.8, 0.4, "December 2021", 16, SECONDARY, align=PP_ALIGN.CENTER)
add_text(s, 6.7, 1.95, 5.9, 0.7, "1   The warning was published.", 24, INK)
add_text(s, 6.7, 2.95, 5.9, 0.7, "2   It reached an agency dashboard.", 24, INK)
add_text(s, 6.7, 3.95, 5.9, 0.7, "3   It never reached her shop floor.", 24, RED, bold=True)
b = add_bar(s, 0.6, 6.4, 12.13, 0.5, ACCENT)
b.text_frame.word_wrap = True
p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "“A warning only helps if people can act.”"
run.font.size = Pt(18); run.font.italic = True; run.font.color.rgb = WHITE; run.font.name = FONT
b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
add_footer(s)
notes(s, "December 2021. Kak Ros runs a kedai runcit in Taman Sri Muda. The government's "
         "river station — two-point-nine kilometres away — crossed WARNING at two in the "
         "afternoon. The warning was published. She lost forty thousand ringgit anyway. Not "
         "because the warning didn't exist — because nobody translated 'river at warning' into "
         "her freezer, her rice sacks, her next three hours. The brief asks us to keep one place "
         "running when the weather turns bad. We picked the flood-prone kedai — and we found why "
         "the warnings already there don't help.")

# ============ SLIDE 3 (system map) ============
s = slide(); add_bg(s)
add_headline(s, "One broken link — the last metre")
# upstream boxes
box_text(s, 0.6, 2.2, 1.2, 0.8, "Rain", 16, WHITE, bold=True, fill=ACCENT_LIGHT, border=None)
box_text(s, 2.0, 2.2, 1.3, 0.8, "Rivers", 16, WHITE, bold=True, fill=ACCENT_LIGHT, border=None)
box_text(s, 3.5, 2.2, 2.6, 0.8, "JPS telemetry\n~200 stations · thresholds", 14, WHITE, bold=True, fill=ACCENT_LIGHT, border=None)
box_text(s, 6.3, 2.2, 2.0, 0.8, "Agency dashboard", 15, WHITE, bold=True, fill=ACCENT_LIGHT, border=None)
# connectors upstream
for x in (1.8, 3.3, 6.1):
    arrow(s, x, 2.6, x + 0.2, 2.6, color=ACCENT_LIGHT, weight=2.0)
# THE LAST METRE red dashed arrow
add_text(s, 8.2, 1.5, 2.4, 0.5, "THE LAST METRE", 15, RED, bold=True, align=PP_ALIGN.CENTER)
arrow(s, 8.35, 2.6, 9.15, 2.6, color=RED, weight=3.5, dash=True)
oval(s, 8.15, 1.95, 2.5, 1.35, color=RED, weight=2.5)
# downstream vertical stack (right)
ds = ["Shop floor", "Owner action", "Damage / Claim", "Loan / Recovery"]
for i, label in enumerate(ds):
    ty = 2.0 + i * 0.9
    box_text(s, 10.9, ty, 2.0, 0.7, label, 15, INK, fill=GREY_BOX, border=CARD_BORDER)
    if i > 0:
        arrow(s, 11.9, ty - 0.2, 11.9, ty, color=SECONDARY, weight=1.75)
# callout
add_text(s, 6.6, 5.85, 6.2, 0.7, "Highest leverage · Lowest cost · Our intervention point",
         18, ACCENT, bold=True, align=PP_ALIGN.CENTER)
# iceberg inset
add_text(s, 0.6, 4.0, 3.4, 0.35, "Iceberg", 12, SECONDARY)
ice = [("Event", ACCENT_LIGHT, WHITE), ("Pattern", ACCENT, WHITE),
       ("Structure", RGBColor(0x03, 0x69, 0x9C), WHITE), ("Mental model", INK, WHITE)]
for i, (lab, fill, col) in enumerate(ice):
    box_text(s, 0.6, 4.4 + i * 0.46, 3.4, 0.4, lab, 13, col, fill=fill, border=None)
add_footer(s)
notes(s, "So we mapped the system. Rain feeds rivers; Malaysia has world-class telemetry — "
         "around two hundred live river stations with published danger thresholds. That data "
         "reaches an agency dashboard — and then it stops. Between the river gauge and the shop "
         "floor is a gap we call the last metre. Everything upstream is solved and expensive. "
         "Everything downstream — a flooded shop, a denied claim, a panic loan — flows from this "
         "one broken link. On the iceberg: the event is a flooded shop; the pattern is the same "
         "shops flooding every monsoon despite warnings; the structure is warnings that are "
         "river-centric with no shop-level translation and no evidence trail; the mental model is "
         "'banjir is fate.' The brief's own question is 'what if floods double?' — they will, and "
         "every extra flood widens exactly this gap. So the last metre is our intervention point: "
         "the highest-leverage, lowest-cost place in the whole system.")

# ============ SLIDE 4 ============
s = slide(); add_bg(s)
add_headline(s, "A flood isn't water — it's money")
chain = [("Weather event", 0.6, 2.6), ("What's hit", 3.5, 2.4), ("Effect", 6.2, 2.0), ("Cost", 8.5, 2.0)]
for label, l, w in chain:
    box_text(s, l, 1.8, w, 0.7, label, 18, WHITE, bold=True, fill=ACCENT, border=None)
for x in (3.3, 5.95, 8.3):
    arrow(s, x, 2.15, x + 0.18, 2.15, color=ACCENT, weight=2.0)
effects = [("Damaged goods", ""), ("Lost sales", "days shut"),
           ("Cleanup costs", ""), ("Costlier loan", "no evidence → claim denied")]
for i, (title, sub) in enumerate(effects):
    l = 0.6 + i * 2.955
    box(s, l, 3.0, 2.85, 1.5)
    add_text(s, l + 0.15, 3.25, 2.55, 0.6, title, 18, INK, bold=True, align=PP_ALIGN.CENTER)
    if sub:
        add_text(s, l + 0.15, 3.85, 2.55, 0.5, sub, 14, RED if "evidence" in sub else SECONDARY, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 5.05, 12.13, 1.0, "RM 6.1 billion", 60, INK, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 6.25, 12.13, 0.5, "national flood losses, Dec 2021 — the market we address",
         16, SECONDARY, align=PP_ALIGN.CENTER)
add_footer(s)
notes(s, "And a flood isn't water, it's money. The chain: weather event, what's hit, effect, "
         "cost. Water reaches her floor-level freezer and stock and triggers all four money "
         "effects at once — damaged goods, lost sales while she's shut, cleanup costs, and a "
         "costlier loan because with no evidence her claim is denied. December 2021 was "
         "six-point-one billion ringgit nationally. Every ringgit started as a shop like hers — "
         "and that RM 6.1 billion is the market we're addressing.")

# ============ SLIDE 5 ============
s = slide(); add_bg(s)
add_headline(s, "AI is never in the storm path")
box(s, 0.6, 1.9, 5.7, 3.2, fill=LIGHT_BLUE, border=None)
add_text(s, 0.9, 2.15, 5.1, 0.5, "DRY DAY — AI thinks", 20, ACCENT, bold=True)
add_text(s, 0.9, 2.8, 5.1, 2.0, "5 photos → asset risk map → owner confirms → cached tiered playbooks",
         18, INK)
box(s, 7.03, 1.9, 5.7, 3.2, fill=LIGHT_GREY, border=None)
add_text(s, 7.33, 2.15, 5.1, 0.5, "STORM DAY — dumb & deterministic", 20, INK, bold=True)
add_text(s, 7.33, 2.8, 5.1, 1.0, "threshold → cache lookup → send", 18, INK)
box_text(s, 3.4, 3.25, 6.5, 0.7, "AI is never in the storm path", 18, WHITE, bold=True, fill=ACCENT, border=None)
box_text(s, 0.6, 5.4, 12.13, 1.0,
         "After the flood: 5 photos → claim report (evidence = river authority telemetry)",
         18, INK, fill=BG_SOFT, align=PP_ALIGN.LEFT)
add_footer(s)
notes(s, "BanjirKawan. The design principle answers the brief's own systems angle — 'the fix "
         "needs power, phones and roads too' — so the AI is never in the storm path. On a dry "
         "day, AI vision turns five phone photos into a risk map of every asset — its height off "
         "the floor, its ringgit value — and the owner confirms it, so nothing is hallucinated. "
         "We pre-compute tiered action playbooks and cache them. When the flood comes, the alert "
         "path is dumb, deterministic code — threshold, cache lookup, send. And after the flood, "
         "five photos become an itemised claim report whose evidence is the river authority's own "
         "telemetry.")

# ============ SLIDE 6 ============
s = slide(); add_bg(s)
add_headline(s, "Not a concept — deployed, on real data")
mp4 = asset("demo.mp4")
if mp4:
    poster = asset("map.png")
    s.shapes.add_movie(mp4, Inches(0.6), Inches(1.8), Inches(7.0), Inches(4.3),
                       poster_frame_image=poster) if poster else \
        s.shapes.add_movie(mp4, Inches(0.6), Inches(1.8), Inches(7.0), Inches(4.3))
else:
    box_text(s, 0.6, 1.8, 7.0, 4.3,
             "[ 15–20s recording:\nSIMULATE FLOOD → phone buzz → checklist → dashboard tick ]",
             20, WHITE, fill=INK, border=None)
add_text(s, 0.6, 6.2, 7.0, 0.4, "301 live river stations · 2 real escalations caught this week",
         15, SECONDARY)
stills = [("Live danger map", "map.png"),
          ("Claim report — evidence = JPS telemetry", "report.png"),
          ("Metrics dashboard", "metrics.png")]
for i, (cap, fn) in enumerate(stills):
    ty = 1.8 + i * 1.45
    img = asset(fn)
    if img:
        s.shapes.add_picture(img, Inches(7.9), Inches(ty), Inches(4.8), Inches(1.15))
    else:
        box_text(s, 7.9, ty, 4.8, 1.15, cap, 14, INK, fill=BG_SOFT)
    add_text(s, 7.9, ty + 1.02, 4.8, 0.3, cap, 12, SECONDARY)
add_footer(s)
notes(s, "And this isn't a concept — we've built and deployed it. It's monitoring 301 real "
         "river stations, it's already caught two real escalations this week, and here's a "
         "warning firing to a shop owner's phone and the action ticking off, live. Built product, "
         "real data — that de-risks everything I'm about to say about the business.")

# ============ SLIDE 7 (business) ============
s = slide(); add_bg(s)
add_headline(s, "B2B2C — and the shop is the smallest payer")
# card 1 anchor
box(s, 0.6, 1.8, 4.0, 3.0, fill=LIGHT_BLUE, border=ACCENT, border_w=2.5)
box_text(s, 0.85, 2.0, 1.2, 0.35, "ANCHOR", 11, WHITE, bold=True, fill=ACCENT, border=None)
add_text(s, 0.85, 2.45, 3.5, 0.5, "Insurers / takaful", 22, ACCENT, bold=True)
add_text(s, 0.85, 3.05, 3.5, 1.6,
         "Completed checklist = claim avoided. Verified report = claim in minutes, not "
         "adjuster-weeks. Priced per policy / per claim avoided.", 15, INK)
# card 2
box(s, 4.75, 1.8, 3.9, 3.0)
add_text(s, 5.0, 2.1, 3.4, 0.5, "State & local councils", 20, INK, bold=True)
add_text(s, 5.0, 2.75, 3.4, 1.5, "District-wide SME resilience at software cost.", 16, INK)
# card 3
box(s, 8.8, 1.8, 3.93, 3.0)
add_text(s, 9.05, 2.1, 3.4, 0.5, "Kedai — freemium", 20, INK, bold=True)
add_text(s, 9.05, 2.7, 3.4, 0.9, "Basic alerts free. Plan + claim features paid.", 16, INK)
add_text(s, 9.05, 3.7, 3.4, 0.6, "RM 19 / month", 22, ACCENT, bold=True)
# economics chips
chips = [("Marginal cost / shop / flood ≈ 0", "one dry-day AI call, then DB lookup + one message"),
         ("Data moat", "site graphs + behavioural completion data nobody else has"),
         ("Go to market", "one takaful + one council pilot, measured vs the next street")]
for i, (t1, t2) in enumerate(chips):
    l = 0.6 + i * 4.075
    box(s, l, 5.2, 3.9, 1.3)
    add_text(s, l + 0.2, 5.35, 3.5, 0.5, t1, 16, INK, bold=True)
    add_text(s, l + 0.2, 5.85, 3.5, 0.6, t2, 13, SECONDARY)
add_footer(s)
notes(s, "The business. It's B2B2C, and the systems view hands us three payers — and the shop "
         "is the smallest of them. One, insurers and takaful — this is the anchor. A completed "
         "checklist is a claim that never happens; a verified, telemetry-backed report is a claim "
         "that costs minutes instead of adjuster-weeks. We sell risk reduction and faster, cheaper "
         "claims — priced per policy or per claim avoided. Two, state and local councils — flood "
         "resilience for their SMEs at software cost, a whole district for a rounding error in a "
         "disaster budget. Three, the kedai, freemium at nineteen ringgit a month. And the "
         "economics are the point: onboarding uses one AI call on a dry day, and every storm-time "
         "alert after that is a database lookup and one message — the marginal cost per shop per "
         "flood is effectively zero. So it scales to the whole six-point-one-billion problem "
         "without cost scaling with it. And it compounds: every shop builds a data moat — site "
         "graphs and behavioural data no competitor and no government dashboard has. We go to "
         "market through one takaful partner and one council pilot — a single flood-prone street, "
         "measured against the street next to it.")

# ============ SLIDE 8 (metrics) ============
s = slide(); add_bg(s)
add_headline(s, "We measure impact, not adjectives")
tiles = [("RM 10,600–22,200", "protected per event"),
         ("3h 12m", "warning lead time"),
         ("2.0s", "dispatch latency"),
         ("83%", "actions completed"),
         ("90s", "to first action"),
         ("4 min", "claim report (vs weeks)")]
for i, (num, lab) in enumerate(tiles):
    col = i % 3; row = i // 3
    l = 0.6 + col * 3.95
    t = 1.9 + row * 1.95
    box(s, l, t, 3.7, 1.7)
    add_text(s, l + 0.15, t + 0.25, 3.4, 0.8, num, 34, ACCENT, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, l + 0.15, t + 1.1, 3.4, 0.4, lab, 16, SECONDARY, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 5.9, 12.13, 0.5,
         "The numbers a takaful actuary and a council underwrite against.", 18, INK, align=PP_ALIGN.CENTER)
add_footer(s)
notes(s, "And we measure impact, not adjectives — these are pulled live from our running "
         "system: RM 10,600 to 22,200 protected in a single event, a three-hour-twelve-minute "
         "warning window, 83% of actions completed, first action in 90 seconds, and a claim report "
         "produced in four minutes versus weeks by hand. Those are the numbers a takaful actuary "
         "and a council both underwrite against.")

# ============ SLIDE 9 (systems) ============
s = slide(); add_bg(s)
add_headline(s, "Built as a system, not an app")
box(s, 0.6, 1.8, 5.8, 1.15, fill=LIGHT_GREEN, border=None)
add_text(s, 0.85, 1.92, 5.4, 0.35, "BALANCING", 14, GREEN, bold=True)
add_text(s, 0.85, 2.3, 5.4, 0.6, "completed checklists → better playbooks (learns each monsoon)", 16, INK)
box(s, 0.6, 3.1, 5.8, 1.15, fill=LIGHT_BLUE, border=None)
add_text(s, 0.85, 3.22, 5.4, 0.35, "REINFORCING", 14, ACCENT, bold=True)
add_text(s, 0.85, 3.6, 5.4, 0.6, "paid claims → trust → more shops (grows from use)", 16, INK)
# worst-day table
rows, cols = 4, 2
tbl_shape = s.shapes.add_table(rows, cols, Inches(6.7), Inches(1.8), Inches(6.0), Inches(2.6))
tbl = tbl_shape.table
tbl.columns[0].width = Inches(2.0); tbl.columns[1].width = Inches(4.0)
data = [("If…", "Then…"),
        ("Data feed dies", "Watchdog: treat heavy rain as a warning"),
        ("Network dies", "Earliest tier fired hours early + paper plan"),
        ("Floods double", "Marginal cost ≈ 0 — plans already cached")]
for r in range(rows):
    for c in range(cols):
        cell = tbl.cell(r, c)
        cell.fill.solid()
        cell.fill.fore_color.rgb = ACCENT if r == 0 else WHITE
        tf = cell.text_frame; tf.word_wrap = True
        p = tf.paragraphs[0]; run = p.add_run(); run.text = data[r][c]
        f = run.font; f.size = Pt(13 if r else 14); f.name = FONT
        f.bold = (r == 0); f.color.rgb = WHITE if r == 0 else INK
b = add_bar(s, 0.6, 5.9, 12.13, 0.6, ACCENT); b.text_frame.vertical_anchor = MSO_ANCHOR.MIDDLE
p = b.text_frame.paragraphs[0]; p.alignment = PP_ALIGN.CENTER
run = p.add_run(); run.text = "3 channels — Telegram · SMS · Paper. No single point of failure."
run.font.size = Pt(18); run.font.bold = True; run.font.color.rgb = WHITE; run.font.name = FONT
add_footer(s)
notes(s, "Three things make this a resilient system, not just an app. Feedback loops: every "
         "completed checklist teaches us which actions owners actually take — a balancing loop "
         "that tightens the playbook each monsoon; every claim we get paid builds the trust that "
         "brings the next shop in — a reinforcing loop. It gets smarter and bigger from use. "
         "Designed for the worst day, not the average: feed dies, a watchdog tells every shop to "
         "treat heavy rain as a warning; network dies, the earliest tier already fired hours "
         "before the water, and there's a laminated plan on the wall. That's the brief's "
         "power-phones-roads point answered with three channels — Telegram, SMS, and paper. No "
         "single point of failure decides whether Kak Ros acts.")

# ============ SLIDE 10 (close) ============
s = slide(); add_bg(s, BG_SOFT)
add_headline(s, "Same monsoon. Same last-metre gap. Same fix.")
pins = ["Shah Alam", "Johor", "Jakarta", "Medan / Bukit Lawang"]
for i, pin in enumerate(pins):
    l = 0.6 + i * 3.03
    box_text(s, l, 2.0, 2.7, 0.8, pin, 18, WHITE, bold=True, fill=ACCENT_LIGHT, border=None)
add_text(s, 0.6, 3.0, 12.13, 0.5,
         "Same monsoon, same last-metre gap, same fix, same three payers.", 15, SECONDARY, align=PP_ALIGN.CENTER)
add_text(s, 1.0, 3.9, 11.33, 1.6,
         "“Every flood warning becomes a survival plan — and every survival plan becomes claim evidence.”",
         30, ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 0.6, 5.8, 12.13, 0.5, "BanjirKawan · Terima kasih", 20, SECONDARY, align=PP_ALIGN.CENTER)
notes(s, "One place. One broken link — the last metre. One intervention, and we've already "
         "built it, deployed it, and put a business around it. Every flood warning becomes a "
         "shop's personal survival plan, and every survival plan becomes its claim evidence. The "
         "rivers are already telling us. BanjirKawan makes sure Kak Ros hears it in time — and "
         "gets paid after. Terima kasih.")

out = os.path.join(ROOT, "docs", "BanjirKawan.pptx")
prs.save(out)
print("Saved:", out, "| slides:", len(prs.slides._sldIdLst))
